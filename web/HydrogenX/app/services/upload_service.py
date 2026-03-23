import csv
import hashlib
import io
import json
import mimetypes
import re
import uuid
from pathlib import Path

from flask import current_app
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader
from werkzeug.datastructures import FileStorage
from werkzeug.utils import secure_filename
from docx import Document

from ..models import UploadedFile


TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".json",
    ".csv",
    ".tsv",
    ".py",
    ".js",
    ".ts",
    ".html",
    ".css",
    ".xml",
    ".yaml",
    ".yml",
    ".ini",
    ".log",
    ".sql",
}


class UploadValidationError(ValueError):
    pass


class UploadService:
    def save_upload(self, file_storage: FileStorage, user_id: int) -> UploadedFile:
        if not file_storage or not getattr(file_storage, "filename", ""):
            raise UploadValidationError("未接收到有效文件。")

        original_name = Path(file_storage.filename.strip()).name
        extension = Path(original_name).suffix.lower()
        safe_name = self.sanitize_filename(original_name, extension=extension)
        allowed_extensions = current_app.config["HYDROGEN_ALLOWED_UPLOAD_EXTENSIONS"]
        if extension not in allowed_extensions:
            allowed_text = "、".join(sorted(ext.lstrip(".") for ext in allowed_extensions))
            raise UploadValidationError(f"暂不支持该文件类型：{extension or '无扩展名'}。支持：{allowed_text}。")

        binary = file_storage.read()
        if not binary:
            raise UploadValidationError("上传文件为空。")

        per_file_limit = int(current_app.config["HYDROGEN_MAX_UPLOAD_SIZE_BYTES"])
        if len(binary) > per_file_limit:
            raise UploadValidationError(
                f"文件 {original_name} 超过大小限制（{self.human_size(per_file_limit)}）。"
            )

        sha256 = hashlib.sha256(binary).hexdigest()
        mime_type = file_storage.mimetype or mimetypes.guess_type(safe_name)[0] or "application/octet-stream"
        storage_name = f"{uuid.uuid4().hex}_{safe_name}"
        relative_path = f"user_{user_id}/{storage_name}"
        target_dir = self.user_upload_dir(user_id)
        target_dir.mkdir(parents=True, exist_ok=True)
        file_path = target_dir / storage_name
        file_path.write_bytes(binary)

        try:
            extracted_text = self.extract_text(file_path=file_path, filename=safe_name, mime_type=mime_type)
        except UploadValidationError:
            if file_path.exists():
                file_path.unlink()
            raise
        except Exception as exc:  # pragma: no cover - defensive path
            if file_path.exists():
                file_path.unlink()
            raise UploadValidationError(f"文件解析失败：{exc}")

        extracted_text = self.normalize_text(extracted_text)
        if not extracted_text:
            if file_path.exists():
                file_path.unlink()
            raise UploadValidationError("未能从文件中提取到可读文本，请上传文本版 PDF/Word/Excel/代码或纯文本文件。")

        max_chars = int(current_app.config["HYDROGEN_MAX_FILE_TEXT_CHARS"])
        extracted_text = extracted_text[:max_chars]

        upload = UploadedFile(
            user_id=user_id,
            filename=safe_name,
            storage_name=storage_name,
            relative_path=relative_path,
            extension=extension or "",
            mime_type=mime_type,
            size_bytes=len(binary),
            sha256=sha256,
            extracted_text=extracted_text,
            extracted_chars=len(extracted_text),
            extraction_status="ready",
            extraction_error=None,
        )
        return upload

    def build_prompt_with_uploads(self, prompt: str, uploads: list[UploadedFile]) -> str:
        prompt = (prompt or "").strip()
        usable_uploads = [upload for upload in uploads if upload and upload.extraction_status == "ready" and upload.extracted_text]
        if not usable_uploads:
            return prompt

        max_total_chars = int(current_app.config["HYDROGEN_MAX_TOTAL_CONTEXT_CHARS"])
        sections = []
        used_chars = 0

        for index, upload in enumerate(usable_uploads, start=1):
            remaining = max_total_chars - used_chars
            if remaining <= 0:
                break

            header = (
                f"[Uploaded File {index}]\n"
                f"Filename: {upload.filename}\n"
                f"MIME Type: {upload.mime_type or 'application/octet-stream'}\n"
                "Extracted Content:\n"
            )
            header_len = len(header)
            if header_len >= remaining:
                break

            body_limit = remaining - header_len
            body = (upload.extracted_text or "")[:body_limit]
            if not body.strip():
                continue

            sections.append(f"{header}{body}\n[/Uploaded File {index}]")
            used_chars += header_len + len(body)

        if not sections:
            return prompt

        joined_sections = "\n\n".join(sections)
        return (
            "以下是用户本轮上传文件中提取出的可读内容。请把这些文件视为本轮问题的上下文依据，优先结合文件内容回答；"
            "若文件证据不足，请明确说明，不要臆造。\n\n"
            "<hydrogenx-user-uploads>\n"
            f"{joined_sections}\n"
            "</hydrogenx-user-uploads>\n\n"
            f"用户问题：\n{prompt}"
        )

    def upload_root(self) -> Path:
        return Path(current_app.root_path).parent / current_app.config["HYDROGEN_UPLOADS_DIR"]

    def user_upload_dir(self, user_id: int) -> Path:
        return self.upload_root() / f"user_{user_id}"

    def upload_abspath(self, upload: UploadedFile) -> Path:
        return self.upload_root() / upload.relative_path

    def delete_upload_file(self, upload: UploadedFile) -> None:
        path = self.upload_abspath(upload)
        if path.exists() and path.is_file():
            path.unlink()

    def sanitize_filename(self, filename: str, extension: str | None = None) -> str:
        original_name = Path(filename or "").name
        suffix = (extension if extension is not None else Path(original_name).suffix).lower()
        stem = Path(original_name).stem

        cleaned_stem = secure_filename(stem)
        if not cleaned_stem:
            cleaned_stem = f"upload_{uuid.uuid4().hex[:12]}"

        candidate = f"{cleaned_stem}{suffix}"
        cleaned_candidate = secure_filename(candidate)

        if cleaned_candidate and Path(cleaned_candidate).suffix.lower() == suffix:
            return cleaned_candidate[:180]

        fallback = f"{cleaned_stem}{suffix}"
        return fallback[:180]

    def extract_text(self, file_path: Path, filename: str, mime_type: str) -> str:
        extension = file_path.suffix.lower()
        if extension in TEXT_EXTENSIONS:
            return self.extract_text_file(file_path)
        if extension == ".pdf":
            return self.extract_pdf(file_path)
        if extension == ".docx":
            return self.extract_docx(file_path)
        if extension == ".xlsx":
            return self.extract_xlsx(file_path)
        if extension == ".pptx":
            return self.extract_pptx(file_path)
        raise UploadValidationError(f"暂不支持解析文件：{filename} ({mime_type})")

    def extract_text_file(self, file_path: Path) -> str:
        binary = file_path.read_bytes()
        encodings = ["utf-8", "utf-8-sig", "gb18030", "gbk", "latin-1"]
        for encoding in encodings:
            try:
                return binary.decode(encoding)
            except UnicodeDecodeError:
                continue
        return binary.decode("utf-8", errors="ignore")

    def extract_pdf(self, file_path: Path) -> str:
        reader = PdfReader(str(file_path))
        parts = []
        for index, page in enumerate(reader.pages, start=1):
            page_text = (page.extract_text() or "").strip()
            if page_text:
                parts.append(f"[Page {index}]\n{page_text}")
        return "\n\n".join(parts)

    def extract_docx(self, file_path: Path) -> str:
        document = Document(str(file_path))
        parts = []
        for paragraph in document.paragraphs:
            text = (paragraph.text or "").strip()
            if text:
                parts.append(text)

        for table_index, table in enumerate(document.tables, start=1):
            rows = []
            for row in table.rows:
                row_values = [self.normalize_cell_text(cell.text) for cell in row.cells]
                if any(row_values):
                    rows.append(" | ".join(row_values))
            if rows:
                parts.append(f"[Table {table_index}]\n" + "\n".join(rows))
        return "\n".join(parts)

    def extract_xlsx(self, file_path: Path) -> str:
        workbook = load_workbook(filename=str(file_path), read_only=True, data_only=True)
        parts = []
        max_rows = int(current_app.config["HYDROGEN_XLSX_MAX_ROWS_PER_SHEET"])
        max_cols = int(current_app.config["HYDROGEN_XLSX_MAX_COLS_PER_ROW"])

        for sheet in workbook.worksheets:
            rows = []
            for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                if row_index > max_rows:
                    rows.append(f"... truncated after {max_rows} rows ...")
                    break
                values = [self.normalize_cell_value(value) for value in row[:max_cols]]
                if any(values):
                    rows.append("\t".join(values))
            if rows:
                parts.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
        return "\n\n".join(parts)

    def extract_pptx(self, file_path: Path) -> str:
        presentation = Presentation(str(file_path))
        parts = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = self.normalize_text(shape.text)
                    if text:
                        texts.append(text)
            if texts:
                parts.append(f"[Slide {slide_index}]\n" + "\n".join(texts))
        return "\n\n".join(parts)

    def normalize_text(self, value: str | None) -> str:
        text = (value or "").replace("\x00", " ")
        text = text.replace("\r\n", "\n").replace("\r", "\n")
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()

    def normalize_cell_text(self, value: str | None) -> str:
        return re.sub(r"\s+", " ", (value or "").strip())

    def normalize_cell_value(self, value) -> str:
        if value is None:
            return ""
        if isinstance(value, (dict, list)):
            return json.dumps(value, ensure_ascii=False)
        return self.normalize_cell_text(str(value))

    def human_size(self, size: int) -> str:
        if size < 1024:
            return f"{size} B"
        if size < 1024 * 1024:
            return f"{size / 1024:.1f} KB"
        return f"{size / (1024 * 1024):.1f} MB"
